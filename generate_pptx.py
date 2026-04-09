import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn  # kept for future XML manipulation
import os

# ── Constants ──────────────────────────────────────────────
BG_DARK  = "1A1A2E"
ACCENT   = "4285F4"
WHITE    = "FFFFFF"
GRAY     = "CCCCCC"
LIGHT_GRAY = "9999AA"
CARD_BG  = "22223A"
GREEN    = "34D39A"

NAME   = "Aleksandr Martiushev"
EMAIL  = "a.martiushev@innopolis.university"
GROUP  = "CSE-03"
GITHUB = "https://github.com/alexzhal1/se-toolkit-hackathon"
URL    = "http://10.93.26.98/"

OUTPUT = os.path.join(os.path.dirname(__file__), "StudyBot_Presentation_v2.pptx")

# ── Helper: solid RGB color ────────────────────────────────
def solid_fill(shape_or_cell, color_hex):
    """Set solid fill on a shape or table cell."""
    shape_or_cell.fill.solid()
    shape_or_cell.fill.fore_color.rgb = RGBColor(*bytes.fromhex(color_hex))


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Set dark background on slide master (applies to ALL slides) ──
bg = prs.slide_masters[0].background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# ── Helper functions ───────────────────────────────────────

def add_textbox(slide, left, top, width, height, text, font_size=24,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*bytes.fromhex(color))
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=20,
                    bullet_color=ACCENT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        rd = p.add_run()
        rd.text = "\u25cf  "
        rd.font.size = Pt(font_size)
        rd.font.color.rgb = RGBColor(*bytes.fromhex(bullet_color))
        rd.font.name = "Calibri"
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(font_size)
        rt.font.color.rgb = RGBColor(*bytes.fromhex(GRAY))
        rt.font.name = "Calibri"
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bytes.fromhex(ACCENT))
    shape.line.fill.background()

def add_card(slide, left, top, width, height, text_runs):
    """
    Add a rounded rectangle card with formatted text.
    text_runs: list of (text, size, color, bold)
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    solid_fill(shape, CARD_BG)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.2)
    tf.margin_bottom = Inches(0.15)

    for i, (txt, sz, clr, bld) in enumerate(text_runs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = RGBColor(*bytes.fromhex(clr))
        r.font.bold = bld
        r.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT
    return shape

def make_qr(path, text):
    img = qrcode.make(text, border=1)
    img = img.resize((180, 180))
    img.save(path)

# ── QR Codes ───────────────────────────────────────────────
qr_dir = os.path.join(os.path.dirname(__file__), "qr_temp")
os.makedirs(qr_dir, exist_ok=True)
qr_gh = os.path.join(qr_dir, "qr_github.png")
qr_pr = os.path.join(qr_dir, "qr_product.png")
make_qr(qr_gh, GITHUB)
make_qr(qr_pr, URL)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_accent_line(slide, Inches(1.5), Inches(1.8), Inches(1.5))
add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(1.2),
            "StudyBot", font_size=60, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
            "AI-Powered Study Assistant", font_size=28, color=ACCENT)
add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.4),
            NAME, font_size=20, color=GRAY)
add_textbox(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.4),
            f"{EMAIL}  |  {GROUP}", font_size=18, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — Context
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
            "Context", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(1.2))

# End-user card
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2), [
    ("End-user", 22, ACCENT, True),
    ("University students who want to efficiently understand study material and retain it long-term.", 17, GRAY, False),
])

# Problem card
add_card(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.2), [
    ("Problem", 22, ACCENT, True),
    ("Students spend hours on lecture notes and textbooks without a tool to explain concepts, generate flashcards, or create self-assessment quizzes.", 17, GRAY, False),
])

# Product idea card
add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5), [
    ("Product idea", 22, ACCENT, True),
    ("StudyBot is an AI assistant that takes study material (text, PDF, DOCX), explains it in simple terms, creates flashcards with spaced repetition, and generates multiple-choice quizzes for self-testing.", 18, WHITE, False),
])

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Implementation
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
            "Implementation", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(1.2))

# ── Tech Stack Section ─────────────────────────────────────
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(0.5),
            "Tech Stack", font_size=24, color=ACCENT, bold=True)

stack_items = [
    ("Backend", "Python 3.12 + FastAPI"),
    ("Database", "PostgreSQL + SQLAlchemy (async) + Alembic"),
    ("AI", "DeepSeek API (OpenAI-compatible)"),
    ("Auth", "JWT (bcrypt + PyJWT)"),
    ("Web Client", "React + Vite + TypeScript"),
    ("Deployment", "Docker Compose + Nginx + VM"),
]

# ── Tech Stack Card ─────────────────────────────────────────
stack_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.2))
solid_fill(stack_shape, CARD_BG)
stack_shape.line.fill.background()
tf = stack_shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
tf.margin_right = Inches(0.2)
tf.margin_bottom = Inches(0.1)
for i, (layer, tech) in enumerate(stack_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    r1 = p.add_run()
    r1.text = f"{layer}:  "
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(*bytes.fromhex(WHITE))
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = tech
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(*bytes.fromhex(GRAY))
    r2.font.name = "Calibri"

# ── V1 Card (left, below stack) ─────────────────────────────
add_card(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.4), [
    ("Version 1", 20, GREEN, True),
    ("Upload study text → AI explanation → Generate flashcards.\nWeb interface with JWT authentication.", 16, GRAY, False),
])

# ── V2 Card (right side \u2014 larger) ──────────────────────────
v2_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.5))
solid_fill(v2_shape, CARD_BG)
v2_shape.line.fill.background()
tf = v2_shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.margin_right = Inches(0.2)
tf.margin_bottom = Inches(0.2)

p = tf.paragraphs[0]
r = p.add_run()
r.text = "Version 2 \u2014 New Features"
r.font.size = Pt(20)
r.font.bold = True
r.font.color.rgb = RGBColor(*bytes.fromhex(ACCENT))
r.font.name = "Calibri"

v2_items = [
    "Quiz generation (10 MCQs, single & multi-select)",
    "Interactive quiz UI with score & explanations",
    "Spaced repetition (SM-2 algorithm)",
    "Review page for due flashcards (quality 0\u20135)",
    "Statistics page (progress tracking)",
    "File upload (.pdf, .docx with text extraction)",
    "Public deployment on VM with Nginx",
]
for item in v2_items:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    rd = p2.add_run()
    rd.text = "\u25b8 "
    rd.font.size = Pt(15)
    rd.font.color.rgb = RGBColor(*bytes.fromhex(ACCENT))
    rd.font.name = "Calibri"
    rt = p2.add_run()
    rt.text = item
    rt.font.size = Pt(15)
    rt.font.color.rgb = RGBColor(*bytes.fromhex(GRAY))
    rt.font.name = "Calibri"

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Demo
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
            "Demo", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(1.2))

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(0.5),
            "Pre-recorded video of Version 2 (\u2264 2 min with voice-over):",
            font_size=20, color=GRAY)

steps = [
    "Open the web app \u2192 home page showing uploaded materials",
    "Upload a PDF file \u2192 material is created automatically",
    "Click \"Explain\" \u2192 AI generates a clear explanation as Markdown",
    "Click \"Flashcards\" \u2192 interactive flip cards appear",
    "Open \"Review\" \u2192 start a spaced repetition session (SM-2)",
    "Click \"Generate Quiz\" \u2192 take a quiz, see your score, review mistakes",
]

for i, step in enumerate(steps):
    y = Inches(2.4) + Inches(0.65) * i
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(*bytes.fromhex(ACCENT))
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(i + 1)
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*bytes.fromhex(WHITE))
    r.font.name = "Calibri"

    add_textbox(slide, Inches(2.0), y, Inches(10), Inches(0.5),
                step, font_size=19, color=GRAY)

# Video placeholder
vp = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(2.2), Inches(3.5), Inches(2.5))
solid_fill(vp, CARD_BG)
vp.line.fill.background()
tf = vp.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.8)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "\U0001F3AC\n[Insert demo\nvideo here]"
r.font.size = Pt(16)
r.font.color.rgb = RGBColor(*bytes.fromhex(LIGHT_GRAY))
r.font.name = "Calibri"

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Links
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
            "Links", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(1.2))

# GitHub card
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5), [
    ("GitHub Repository", 22, ACCENT, True),
    (GITHUB, 16, GRAY, False),
])
slide.shapes.add_picture(qr_gh, Inches(4.8), Inches(3.2), Inches(1.4), Inches(1.4))

# Product card
add_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.5), [
    ("Deployed Product (V2)", 22, GREEN, True),
    (URL, 16, GRAY, False),
])
slide.shapes.add_picture(qr_pr, Inches(11.0), Inches(3.2), Inches(1.4), Inches(1.4))

# ── Save ───────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved \u2192 {OUTPUT}")

# cleanup
import shutil
shutil.rmtree(qr_dir)
